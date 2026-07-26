"""PowerPoint text extraction with slide structure preserved.

The old extractor flattened every text shape into one string.  That remains
available through ``extract_slide_texts`` for backward compatibility, while
``extract_slide_structures`` exposes the title and body groups needed by the
slide-aware chunker.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER


TITLE_PLACEHOLDER_TYPES = {
    PP_PLACEHOLDER.TITLE,
    PP_PLACEHOLDER.CENTER_TITLE,
    PP_PLACEHOLDER.VERTICAL_TITLE,
}


def _paragraph_lines(shape) -> list[str]:
    """Return non-empty paragraphs while preserving line/bullet boundaries."""
    lines: list[str] = []
    for paragraph in shape.text_frame.paragraphs:
        line = "".join(run.text for run in paragraph.runs).strip()
        if line:
            lines.append(line)
    return lines


def _is_title_placeholder(shape) -> bool:
    if not getattr(shape, "is_placeholder", False):
        return False
    try:
        return shape.placeholder_format.type in TITLE_PLACEHOLDER_TYPES
    except (AttributeError, KeyError, ValueError):
        return False


def _shape_position(shape) -> dict[str, int]:
    return {
        "left": int(getattr(shape, "left", 0)),
        "top": int(getattr(shape, "top", 0)),
        "width": int(getattr(shape, "width", 0)),
        "height": int(getattr(shape, "height", 0)),
    }


def _iter_leaf_shapes(shapes):
    """Yield text/table shapes recursively, including shapes inside PPT groups."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_leaf_shapes(shape.shapes)
        else:
            yield shape


def _table_lines(shape) -> list[str]:
    lines: list[str] = []
    for row in shape.table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        if any(cells):
            lines.append(" | ".join(cells))
    return lines


def _fallback_title(elements: list[dict[str, Any]], slide_height: int) -> str:
    """Choose a short text box near the top when the slide has no title placeholder."""
    candidates = [
        element
        for element in elements
        if element["type"] == "text"
        and element["top"] <= slide_height * 0.28
        and len(element["text"]) <= 120
    ]
    if not candidates:
        return ""
    candidates.sort(key=lambda element: (element["top"], element["left"]))
    return candidates[0]["text"]


def extract_slide_structures(pptx_path: str | Path) -> list[dict[str, Any]]:
    """Extract title, body groups, tables, and layout coordinates per slide."""
    path = Path(pptx_path)
    if not path.exists():
        raise FileNotFoundError(f"PPTX file not found: {path}")

    presentation = Presentation(str(path))
    slides: list[dict[str, Any]] = []

    for slide_no, slide in enumerate(presentation.slides, start=1):
        title = ""
        elements: list[dict[str, Any]] = []

        for shape in _iter_leaf_shapes(slide.shapes):
            position = _shape_position(shape)

            if getattr(shape, "has_table", False):
                lines = _table_lines(shape)
                if lines:
                    elements.append(
                        {
                            "shape_id": int(shape.shape_id),
                            "type": "table",
                            "lines": lines,
                            "text": "\n".join(lines),
                            **position,
                        }
                    )
                continue

            if not getattr(shape, "has_text_frame", False):
                continue

            lines = _paragraph_lines(shape)
            if not lines:
                continue

            text = "\n".join(lines)
            if _is_title_placeholder(shape):
                title = text
                continue

            elements.append(
                {
                    "shape_id": int(shape.shape_id),
                    "type": "text",
                    "lines": lines,
                    "text": text,
                    **position,
                }
            )

        elements.sort(key=lambda element: (element["top"], element["left"]))

        if not title:
            title = _fallback_title(elements, int(presentation.slide_height))
            if title:
                # Do not duplicate a heuristic title in both title and body.
                for index, element in enumerate(elements):
                    if element["type"] == "text" and element["text"] == title:
                        elements.pop(index)
                        break

        flat_parts = [part for part in [title, *[e["text"] for e in elements]] if part]
        slides.append(
            {
                "slide_no": slide_no,
                "title": title,
                "body_groups": elements,
                "text": "\n".join(flat_parts),
            }
        )

    return slides


def extract_slide_texts(pptx_path: str | Path) -> list[dict[str, Any]]:
    """Original flat extractor used by the baseline chunking strategies.

    This keeps the previous project's behavior: text and tables are
    concatenated per slide without title/body classification.
    """
    path = Path(pptx_path)
    if not path.exists():
        raise FileNotFoundError(f"PPTX file not found: {path}")

    presentation = Presentation(str(path))
    results: list[dict[str, Any]] = []

    for slide_no, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                parts.extend(_paragraph_lines(shape))
            elif getattr(shape, "has_table", False):
                parts.extend(_table_lines(shape))

        results.append({"slide_no": slide_no, "text": "\n".join(parts)})

    return results


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="Extract structured text from a PPTX file.")
    parser.add_argument("pptx", help="Path to the PPTX file")
    args = parser.parse_args()

    extracted = extract_slide_structures(args.pptx)
    print(f"Extracted {len(extracted)} slides")
    for sample in extracted[:3]:
        print(
            f"slide={sample['slide_no']} title={sample['title']!r} "
            f"groups={len(sample['body_groups'])} chars={len(sample['text'])}"
        )
